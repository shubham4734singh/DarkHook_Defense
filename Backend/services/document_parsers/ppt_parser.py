# ================================================================
# ppt_parser.py — DarkHOOK_ Defence
# Version  : 2.0 — Enterprise Grade
# Purpose  : PowerPoint file phishing and malware detection
#            using 14 industry-standard techniques
#
# Technique 1  -> File Type Validation
# Technique 2  -> Metadata Analysis
# Technique 3  -> Macro Detection
# Technique 4  -> VBA Behavior Analysis
# Technique 5  -> Animation Trigger Detection (PPT UNIQUE)
# Technique 6  -> Embedded Object Detection
# Technique 7  -> External Link Analysis
# Technique 8  -> Phishing Content Detection
# Technique 9  -> URL Detection in Slides
# Technique 10 -> Hidden Slide Detection (PPT UNIQUE)
# Technique 11 -> Obfuscation Detection
# Technique 12 -> Action Button Analysis (PPT UNIQUE)
# Technique 13 -> Media File Analysis
# Technique 14 -> Attack Chain Inference + Scoring
#
# Libraries: python-pptx, oletools, zipfile,
#            hashlib, re, math
# ================================================================


# ----------------------------------------------------------------
# IMPORTS
# ----------------------------------------------------------------

import re
import os
import math
import zipfile
import hashlib
from collections import Counter
from urllib.parse import urlparse

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    from oletools.olevba import VBA_Parser
    OLETOOLS_AVAILABLE = True
except ImportError:
    OLETOOLS_AVAILABLE = False

try:
    from .deobfuscator import run_deobfuscation_pipeline
    DEOBFUSCATOR_AVAILABLE = True
except ImportError:
    try:
        from services.document_parsers.deobfuscator import run_deobfuscation_pipeline
        DEOBFUSCATOR_AVAILABLE = True
    except ImportError:
        DEOBFUSCATOR_AVAILABLE = False


# ================================================================
# CONFIGURATION
# ================================================================
# NOTE: Scoring weights are now centralized in scorer.py
WEIGHTS = {
    # Kept for documentation only
    "pps_file": 20,
}


# ----------------------------------------------------------------
# PHISHING KEYWORDS
# ----------------------------------------------------------------

PHISHING_KEYWORDS = {

    "account_threats": [
        "verify your account",
        "confirm your account",
        "account suspended",
        "account will be closed",
        "account has been compromised",
        "account limited",
        "reactivate your account",
        "unusual activity detected",
    ],

    "urgency_phrases": [
        "urgent action required",
        "immediate action required",
        "act now",
        "respond immediately",
        "expires today",
        "final warning",
        "last chance",
        "do not ignore",
        "within 24 hours",
        "time sensitive",
    ],

    "credential_harvesting": [
        "enter your password",
        "confirm your password",
        "reset your password",
        "password expired",
        "enter your credentials",
        "verify your identity",
        "sign in to continue",
        "enter your username",
        "enter your email",
    ],

    "financial_terms": [
        "bank account details",
        "credit card details",
        "debit card number",
        "payment details required",
        "update payment method",
        "transaction failed",
        "wire transfer",
        "gift card",
        "bitcoin payment",
        "refund pending",
    ],

    "macro_lure": [
        "enable content",
        "enable macros",
        "click enable",
        "enable editing",
        "content is disabled",
        "enable to view",
        "protected document",
        "restricted content",
        "click enable content",
    ],

    "india_specific": [
        "aadhar number",
        "pan card details",
        "kyc verification",
        "kyc update required",
        "upi details",
        "enter otp",
        "otp verification",
        "neft",
        "imps",
    ],

    "legal_threats": [
        "legal action will be taken",
        "police complaint filed",
        "court notice",
        "government notice",
        "income tax department",
        "tax refund",
        "warrant issued",
    ],

    "reward_tricks": [
        "you have won",
        "prize money",
        "claim your reward",
        "lottery winner",
        "congratulations you won",
        "cash prize",
        "free gift",
    ],
}


# Suspicious author names
SUSPICIOUS_AUTHORS = [
    "admin", "user", "test", "unknown",
    "administrator", "root", "guest",
    "default", "temp", "anon",
]


# Dangerous VBA API calls
DANGEROUS_VBA_APIS = [
    "shell", "wscript.shell", "createobject",
    "powershell", "urldownloadtofile",
    "winexec", "cmd.exe", "shellexecute",
    "environ", "winhttprequest", "xmlhttp",
    "getobject", "sendkeys", "appactivate",
    "vbscript", "cscript", "wscript",
    "regwrite", "regread", "filesystemobject",
    "openastext", "writeline", "saveas",
]


# PPT specific auto-run macro names
PPT_AUTOOPEN_NAMES = [
    "auto_open", "autoopen",
    "presentation_open", "slideshowbegin",
    "onslideshow", "auto_run",
    "slideshow_open", "workbook_open",
    "document_open",
]


# URL shorteners
URL_SHORTENERS = [
    "bit.ly", "tinyurl.com", "t.co",
    "goo.gl", "ow.ly", "is.gd",
    "buff.ly", "rebrand.ly", "cutt.ly",
    "shorturl.at", "tiny.cc", "rb.gy",
    "qrfy.io", "qrfy.com", "qr.io",
    "qrco.de", "qrd.by",
]


# Suspicious TLDs
SUSPICIOUS_TLDS = [
    ".xyz", ".top", ".ru", ".tk",
    ".ml", ".ga", ".cf", ".gq",
    ".pw", ".click", ".download",
    ".loan", ".work", ".party",
]


# Safe external domains
SAFE_DOMAINS = [
    "schemas.openxmlformats.org",
    "schemas.microsoft.com",
    "purl.org",
    "www.w3.org",
    "microsoft.com",
    "office.com",
    "openxmlformats.org",
]


# Media file extensions
MEDIA_EXTENSIONS = [
    ".png", ".jpg", ".jpeg", ".gif",
    ".bmp", ".tiff", ".webp",
    ".mp4", ".avi", ".mov", ".wmv",
    ".mp3", ".wav", ".wma",
    ".svg", ".emf", ".wmf",
]


# ================================================================
# HELPER FUNCTIONS
# ================================================================

def calculate_entropy(data):
    """
    Shannon entropy — measures randomness.
    High entropy = possible encoded/encrypted data.
    """
    if not data or len(data) < 20:
        return 0.0
    counter = Counter(data)
    length  = len(data)
    entropy = -sum(
        (c / length) * math.log2(c / length)
        for c in counter.values()
    )
    return round(entropy, 2)


def is_ip_url(url):
    """Returns True if URL uses IP address instead of domain."""
    try:
        parsed = urlparse(url if "://" in url else "http://" + url)
        host = parsed.hostname or ""
        return bool(re.match(
            r'^\d{1,3}\.\d{1,3}\.\d{1,3}\.\d{1,3}$',
            host
        ))
    except Exception:
        return False


def is_safe_domain(url):
    """Returns True if URL belongs to safe/trusted domain."""
    try:
        parsed = urlparse(url if "://" in url else "http://" + url)
        domain = (parsed.hostname or "").lower()
        if not domain:
            return False
        for safe in SAFE_DOMAINS:
            if domain == safe or domain.endswith("." + safe):
                return True
        return False
    except Exception:
        return False


def analyze_url(url):
    """Full URL threat analysis."""
    url_findings = []
    url_details  = []

    try:
        parsed = urlparse(url if "://" in url else "http://" + url)
        domain = (parsed.hostname or parsed.netloc or "").lower()
        path   = (parsed.path or "").lower()

        if is_ip_url(url):
            url_findings.append("ip_based_url")
            url_details.append("IP-based URL: " + url)

        for s in URL_SHORTENERS:
            if s in domain:
                url_findings.append("shortened_url")
                url_details.append(
                    "URL shortener (" + s + "): " + url
                )
                break

        for tld in SUSPICIOUS_TLDS:
            if domain.endswith(tld):
                url_findings.append("suspicious_tld")
                url_details.append(
                    "Suspicious TLD (" + tld + "): " + url
                )
                break

        if "@" in url:
            url_findings.append("at_symbol_trick")
            url_details.append("@ trick: " + url)

        if url.startswith("http://"):
            url_findings.append("suspicious_url")
            url_details.append("Insecure HTTP: " + url)

        if len(url) > 200:
            url_findings.append("suspicious_url")
            url_details.append(
                "Very long URL: " + url[:60] + "..."
            )

        fake_words = [
            "login", "signin", "verify", "secure",
            "account", "update", "confirm", "banking",
            "password", "credential",
        ]
        for word in fake_words:
            if word in path:
                url_findings.append("suspicious_url")
                url_details.append(
                    "Login keyword in URL: " + url
                )
                break

    except Exception as e:
        url_details.append("URL error: " + str(e))

    return url_findings, url_details


# ================================================================
# TECHNIQUE 1 — File Type Validation
# ================================================================

def technique1_file_validation(file_path):
    """
    Validates PowerPoint file type and structure.
    Checks extension, ZIP structure and file signature.
    """
    findings = []
    details  = []

    details.append("--- TECHNIQUE 1: FILE VALIDATION ---")

    filename = os.path.basename(file_path)
    ext      = os.path.splitext(filename)[1].lower()

    valid_extensions = [
        ".pptx", ".ppt", ".pps", ".ppsx",
        ".pptm", ".ppsm", ".potx", ".potm",
    ]

    if ext not in valid_extensions:
        findings.append("invalid_ppt_format")
        details.append(
            "Invalid PowerPoint extension: " + ext
        )
    else:
        details.append("Extension valid: " + ext)

    # Flag macro-enabled formats
    if ext in [".pptm", ".ppsm", ".potm"]:
        findings.append("vba_macro_detected")
        details.append(
            "Macro-enabled format: " + ext +
            " — macros may be present!"
        )

    # Flag .pps (auto-runs slideshow)
    if ext in [".pps", ".ppsx"]:
        findings.append("pps_file")
        details.append(
            "PPS format — auto-runs slideshow on open! "
            "Common phishing trick!"
        )

    # Double extension check
    name_no_ext = os.path.splitext(filename)[0]
    if "." in name_no_ext:
        part = name_no_ext.split(".")[-1].strip().lower()
        dangerous = [
            "exe", "dll", "bat", "cmd",
            "ps1", "vbs", "js", "hta",
            "scr", "zip", "rar", "7z", "iso", "pdf",
        ]
        if part in dangerous:
            findings.append("double_extension")
            details.append(
                "Double extension: " + filename
            )
        elif part.isdigit():
            details.append(
                "Version number in name — safe ✅"
            )

    # Check file signatures
    if ext in [".pptx", ".ppsx", ".pptm", ".ppsm", ".potx", ".potm"]:
        try:
            if zipfile.is_zipfile(file_path):
                details.append(
                    "ZIP structure valid ✅"
                )
                with zipfile.ZipFile(
                    file_path, "r"
                ) as z:
                    names = z.namelist()
                    has_slides = any(
                        "ppt/slides/slide" in n or "ppt/presentation.xml" in n
                        for n in names
                    )
                    if has_slides:
                        details.append(
                            "Slide structure found ✅"
                        )
                    else:
                        findings.append(
                            "malformed_zip"
                        )
                        details.append(
                            "No slides found in ZIP — "
                            "suspicious structure!"
                        )
            else:
                findings.append("file_type_mismatch")
                details.append(
                    "File is not valid ZIP — "
                    "possible disguised malware!"
                )
        except Exception as e:
            findings.append("malformed_zip")
            details.append("ZIP error: " + str(e))
    elif ext in [".ppt", ".pps"]:
        try:
            with open(file_path, "rb") as f:
                header = f.read(8)
                if header[:2] == b"MZ":
                    findings.append("file_type_mismatch")
                    findings.append("mz_header_found")
                    details.append("Disguised executable: MZ header found in .ppt file!")
                elif header[:4] != b"\xd0\xcf\x11\xe0":
                    findings.append("file_type_mismatch")
                    details.append("File does not match OLE (.ppt) signature!")
                else:
                    details.append("Legacy OLE signature valid ✅")
        except Exception as e:
            details.append("Header check error: " + str(e))

    details.append(
        "Technique 1 findings: " + str(len(findings))
    )
    return findings, details


# ================================================================
# TECHNIQUE 2 — Metadata Analysis
# ================================================================

def technique2_metadata(file_path):
    """
    Extracts and analyzes PowerPoint metadata.
    Suspicious metadata can reveal attacker tools/identity.
    """
    findings = []
    details  = []

    details.append("")
    details.append("--- TECHNIQUE 2: METADATA ANALYSIS ---")

    ext = os.path.splitext(file_path)[1].lower()

    if ext not in [".pptx", ".ppsx", ".pptm", ".ppsm", ".potx", ".potm"]:
        details.append(
            "Old format — metadata check skipped"
        )
        return findings, details

    try:
        if not zipfile.is_zipfile(file_path):
            details.append("Invalid ZIP — skipping")
            return findings, details

        with zipfile.ZipFile(file_path, "r") as z:
            names = z.namelist()

            # Read core.xml for author/dates
            if "docProps/core.xml" in names:
                core_xml = z.read(
                    "docProps/core.xml"
                ).decode("utf-8", errors="ignore")

                # Extract author
                author_match = re.search(
                    r'<dc:creator[^>]*>(.*?)</dc:creator>',
                    core_xml,
                    re.DOTALL
                )
                if author_match:
                    author = author_match.group(1).strip()
                    details.append(
                        "Author: " + (author or "EMPTY")
                    )
                    if author.lower() in SUSPICIOUS_AUTHORS:
                        findings.append("suspicious_author")
                        details.append(
                            "Suspicious author name: '" +
                            author + "'"
                        )
                    if not author:
                        findings.append("wiped_metadata")
                        details.append(
                            "Author field is empty — wiped!"
                        )
                else:
                    findings.append("wiped_metadata")
                    details.append(
                        "No author found — metadata wiped!"
                    )

                # Extract revision
                rev_match = re.search(
                    r'<cp:revision[^>]*>(.*?)</cp:revision>',
                    core_xml,
                    re.DOTALL
                )
                if rev_match:
                    revision = rev_match.group(1).strip()
                    details.append(
                        "Revision: " + revision
                    )
                    if revision == "1":
                        findings.append(
                            "low_revision_count"
                        )
                        details.append(
                            "Revision = 1 — freshly created, "
                            "never edited!"
                        )

                # Extract dates
                created_match = re.search(
                    r'<dcterms:created[^>]*>(.*?)</dcterms:created>',
                    core_xml,
                    re.DOTALL
                )
                modified_match = re.search(
                    r'<dcterms:modified[^>]*>(.*?)</dcterms:modified>',
                    core_xml,
                    re.DOTALL
                )
                if created_match:
                    details.append(
                        "Created : " +
                        created_match.group(1)
                    )
                if modified_match:
                    details.append(
                        "Modified: " +
                        modified_match.group(1)
                    )

                # Check if created = modified
                if (created_match and modified_match and
                        created_match.group(1).strip() ==
                        modified_match.group(1).strip()):
                    findings.append("metadata_mismatch")
                    details.append(
                        "Created = Modified timestamp — "
                        "file never edited after creation!"
                    )

            else:
                findings.append("wiped_metadata")
                details.append(
                    "core.xml missing — metadata wiped!"
                )

            # Read app.xml for application info
            if "docProps/app.xml" in names:
                app_xml = z.read(
                    "docProps/app.xml"
                ).decode("utf-8", errors="ignore")

                app_match = re.search(
                    r'<Application[^>]*>(.*?)</Application>',
                    app_xml,
                    re.DOTALL
                )
                if app_match:
                    app = app_match.group(1).strip()
                    details.append(
                        "Application: " + app
                    )
                    suspicious_apps = [
                        "unknown",
                    ]
                    for sus_app in suspicious_apps:
                        if sus_app in app.lower():
                            findings.append(
                                "suspicious_metadata"
                            )
                            details.append(
                                "Non-Microsoft app: " + app
                            )
                            break

    except Exception as e:
        details.append("Metadata error: " + str(e))

    details.append(
        "Technique 2 findings: " + str(len(findings))
    )
    return findings, details


# ================================================================
# TECHNIQUE 3 — Macro Detection
# ================================================================

def technique3_macro_detection(file_path):
    """
    Detects VBA macros in PowerPoint file.
    PowerPoint specific auto-run names checked.
    """
    findings    = []
    details     = []
    vba_content = ""

    details.append("")
    details.append(
        "--- TECHNIQUE 3: MACRO DETECTION ---"
    )

    macro_detected = False

    # Check for vbaProject.bin in ZIP
    ext = os.path.splitext(file_path)[1].lower()
    if ext in [".pptx", ".ppsx", ".pptm", ".ppsm", ".potx", ".potm"]:
        try:
            if zipfile.is_zipfile(file_path):
                with zipfile.ZipFile(
                    file_path, "r"
                ) as z:
                    names = z.namelist()
                    if "ppt/vbaProject.bin" in names:
                        macro_detected = True
                        details.append(
                            "vbaProject.bin found — "
                            "VBA macros present!"
                        )
                        vba_content = z.read(
                            "ppt/vbaProject.bin"
                        ).decode(
                            "utf-8", errors="ignore"
                        )
                    else:
                        details.append(
                            "No vbaProject.bin — "
                            "no VBA macros in ZIP ✅"
                        )
        except Exception as e:
            details.append("ZIP error: " + str(e))

    # Use oletools for deep macro analysis
    if OLETOOLS_AVAILABLE:
        try:
            vba_parser = VBA_Parser(file_path)
            try:
                if vba_parser.detect_vba_macros():
                    macro_detected = True
                    details.append(
                        "oletools: VBA macros confirmed!"
                    )

                    extracted_code = ""
                    for (
                        filename, stream_path,
                        vba_filename, vba_code
                    ) in vba_parser.extract_macros():
                        details.append(
                            "Macro in: " + str(vba_filename)
                        )
                        extracted_code += vba_code + "\n"

                        # Check for PPT auto-run names
                        code_lower = vba_code.lower()
                        for name in PPT_AUTOOPEN_NAMES:
                            if name in code_lower and "ppt_autoopen" not in findings:
                                findings.append("ppt_autoopen")
                                details.append(
                                    "CRITICAL: Auto-run macro: '" +
                                    name + "'"
                                )

                    if extracted_code:
                        vba_content = extracted_code

                else:
                    details.append(
                        "oletools: No macros detected ✅"
                    )
            finally:
                vba_parser.close()

        except Exception as e:
            details.append(
                "oletools error: " + str(e)
            )
    else:
        details.append(
            "oletools not available — "
            "install with: pip install oletools"
        )

    if macro_detected:
        findings.append("vba_macro_detected")

    details.append(
        "Technique 3 findings: " + str(len(findings))
    )
    return findings, details, vba_content


# ================================================================
# TECHNIQUE 4 — VBA Behavior Analysis
# ================================================================

def technique4_vba_behavior(vba_content):
    """
    Analyzes VBA macro code for dangerous behavior.
    Checks for shell commands, network calls, etc.
    """
    findings = []
    details  = []

    details.append("")
    details.append(
        "--- TECHNIQUE 4: VBA BEHAVIOR ANALYSIS ---"
    )

    if not vba_content:
        details.append(
            "No VBA content to analyze ✅"
        )
        return findings, details

    code_lower = vba_content.lower()

    # Check dangerous API calls
    for api in DANGEROUS_VBA_APIS:
        if api in code_lower:
            findings.append("suspicious_vba_api")
            details.append(
                "Dangerous API: '" + api + "'"
            )

    # PowerShell detection
    if "powershell" in code_lower:
        findings.append("powershell_in_vba")
        details.append(
            "CRITICAL: PowerShell in VBA macro!"
        )

    # Network call detection
    network_keywords = [
        "xmlhttp", "winhttprequest",
        "urldownloadtofile", "winhttp",
        "internetopen", "httpget",
        "http://", "https://",
    ]
    for keyword in network_keywords:
        if keyword in code_lower:
            findings.append("network_call_in_vba")
            details.append(
                "Network call in macro: '" +
                keyword + "'"
            )
            break

    # File system access
    file_keywords = [
        "open ", "write ", "close ",
        "kill ", "filesystemobject",
        "createtextfile", "environ(",
    ]
    for keyword in file_keywords:
        if keyword in code_lower:
            findings.append("file_system_access")
            details.append(
                "File system access: '" + keyword + "'"
            )
            break

    # Registry access
    registry_keywords = [
        "regwrite", "regread",
        "registry", "hkey_",
        "regedit",
    ]
    for keyword in registry_keywords:
        if keyword in code_lower:
            findings.append("registry_access")
            details.append(
                "Registry access: '" + keyword + "'"
            )
            break

    # Process creation
    process_keywords = [
        "shell(", "shellexecute",
        "winexec", "createobject",
        "wscript.shell",
    ]
    for keyword in process_keywords:
        if keyword in code_lower:
            findings.append("process_creation")
            details.append(
                "Process creation: '" + keyword + "'"
            )
            break

    # Shell command
    if "shell" in code_lower and (
        "cmd" in code_lower or
        "exe" in code_lower
    ):
        findings.append("shell_command")
        details.append(
            "CRITICAL: Shell command with EXE/CMD!"
        )

    # Static De-Obfuscation Pipeline
    if DEOBFUSCATOR_AVAILABLE and vba_content:
        try:
            deob_res = run_deobfuscation_pipeline(vba_content)
            if deob_res.get("deobfuscation_applied"):
                findings.extend(deob_res.get("new_findings", []))
                details.extend(deob_res.get("details", []))
        except Exception as e:
            details.append("Deobfuscator error: " + str(e))

    if not findings:
        details.append(
            "No dangerous VBA behavior detected ✅"
        )

    details.append(
        "Technique 4 findings: " + str(len(findings))
    )
    return findings, details


# ================================================================
# TECHNIQUE 5 — Animation Trigger Detection (PPT UNIQUE)
# ================================================================

def technique5_animation_triggers(file_path):
    """
    Detects malicious animation triggers in slides.
    This technique is UNIQUE to PowerPoint files!

    Attackers hide macro triggers inside animations:
    -> OnMouseOver triggers
    -> Zero-delay auto triggers
    -> Command triggers on slide load
    """
    findings = []
    details  = []

    details.append("")
    details.append(
        "--- TECHNIQUE 5: ANIMATION TRIGGER DETECTION ---"
    )

    ext = os.path.splitext(file_path)[1].lower()
    if ext not in [
        ".pptx", ".ppsx", ".pptm", ".ppsm", ".potx", ".potm",
    ]:
        details.append(
            "Old format — animation check skipped"
        )
        return findings, details

    try:
        if not zipfile.is_zipfile(file_path):
            details.append("Invalid ZIP — skipping")
            return findings, details

        with zipfile.ZipFile(file_path, "r") as z:
            names = z.namelist()

            slide_files = [
                n for n in names
                if re.match(
                    r'ppt/slides/slide\d+\.xml', n
                )
            ]

            details.append(
                "Checking " + str(len(slide_files)) +
                " slides for animation triggers"
            )

            for slide_file in slide_files:
                try:
                    slide_xml = z.read(
                        slide_file
                    ).decode("utf-8", errors="ignore")

                    slide_num = re.search(
                        r'slide(\d+)\.xml', slide_file
                    )
                    snum = (
                        slide_num.group(1)
                        if slide_num else "?"
                    )

                    # Check for timing/animation elements
                    if "<p:timing>" in slide_xml:
                        details.append(
                            "Slide " + snum +
                            ": timing/animation found"
                        )
                        has_suspicious_anim = False

                        # Check for cmd trigger
                        if "<p:cmd" in slide_xml:
                            findings.append(
                                "cmd_trigger_found"
                            )
                            details.append(
                                "CRITICAL: Slide " + snum +
                                " has command trigger!"
                            )
                            has_suspicious_anim = True

                        # Check zero delay trigger on command / action / media nodes
                        has_auto_action = (
                            "<p:cmd" in slide_xml or
                            "ppaction://" in slide_xml.lower() or
                            "<p:cMediaNode" in slide_xml
                        )
                        zero_delay = re.search(
                            r'<p:cTn[^>]*delay="0"',
                            slide_xml
                        )
                        if zero_delay and has_auto_action:
                            findings.append(
                                "zero_delay_trigger"
                            )
                            details.append(
                                "Slide " + snum +
                                ": zero-delay action trigger found!"
                            )
                            has_suspicious_anim = True

                        # Check mouseover trigger in timing
                        if (
                            "mouseover" in
                            slide_xml.lower() or
                            "onmouseover" in
                            slide_xml.lower()
                        ):
                            findings.append(
                                "mouseover_trigger"
                            )
                            details.append(
                                "Slide " + snum +
                                ": mouseover trigger found!"
                            )
                            has_suspicious_anim = True

                        # Generic suspicious animation if risky triggers present
                        if has_suspicious_anim:
                            findings.append(
                                "suspicious_animation"
                            )
                            details.append(
                                "Slide " + snum +
                                ": suspicious animation present"
                            )

                except Exception as e:
                    details.append(
                        "Slide error: " + str(e)
                    )

    except Exception as e:
        details.append(
            "Animation detection error: " + str(e)
        )

    if not findings:
        details.append(
            "No suspicious animation triggers ✅"
        )

    details.append(
        "Technique 5 findings: " + str(len(findings))
    )
    return findings, details


# ================================================================
# TECHNIQUE 6 — Embedded Object Detection
# ================================================================

def technique6_embedded_objects(file_path):
    """
    Detects malicious embedded objects in presentation.
    Checks for EXE files, scripts, OLE objects.
    """
    findings = []
    details  = []

    details.append("")
    details.append(
        "--- TECHNIQUE 6: EMBEDDED OBJECT DETECTION ---"
    )

    ext = os.path.splitext(file_path)[1].lower()
    if ext not in [
        ".pptx", ".ppsx", ".pptm", ".ppsm", ".potx", ".potm",
    ]:
        details.append(
            "Old format — embedded check limited"
        )

    try:
        if not zipfile.is_zipfile(file_path):
            details.append("Invalid ZIP — skipping")
            return findings, details

        with zipfile.ZipFile(file_path, "r") as z:
            names = z.namelist()

            # Check embeddings folder
            embedded = [
                n for n in names
                if "embeddings" in n.lower()
            ]

            details.append(
                "Embedded objects found: " +
                str(len(embedded))
            )

            has_ole = False
            for emb in embedded:
                emb_lower = emb.lower()
                details.append(
                    "Embedded: " + emb
                )

                # Check for OLE objects
                if emb_lower.endswith(".bin"):
                    has_ole = True
                    details.append(
                        "OLE object: " + emb
                    )

                # Check for dangerous file types
                dangerous_embedded = [
                    ".exe", ".dll", ".bat",
                    ".cmd", ".ps1", ".vbs",
                    ".js", ".hta", ".msi",
                ]
                for danger in dangerous_embedded:
                    if emb_lower.endswith(danger):
                        findings.append(
                            "embedded_executable"
                        )
                        details.append(
                            "CRITICAL: Embedded " +
                            danger + ": " + emb
                        )

                # Read and check for MZ header (EXE)
                try:
                    content = z.read(emb)
                    if content[:2] == b"MZ":
                        findings.append(
                            "mz_header_found"
                        )
                        findings.append(
                            "embedded_executable"
                        )
                        details.append(
                            "CRITICAL: MZ header (EXE) "
                            "found in: " + emb
                        )
                except Exception:
                    pass

            # Check for package objects and OLE objects in slides
            has_package = False
            for slide_file in names:
                if not re.match(
                    r'ppt/slides/slide\d+\.xml',
                    slide_file
                ):
                    continue
                try:
                    slide_xml = z.read(
                        slide_file
                    ).decode("utf-8", errors="ignore")

                    if (
                        "<p:oleObj" in slide_xml or
                        "<oleObject" in slide_xml
                    ):
                        has_ole = True
                        details.append(
                            "OLE object in: " + slide_file
                        )

                    # Check for actual OLE package objects (progId="Package" or Packager)
                    if re.search(
                        r'progId="(?i:Package|Packager|Shell)"',
                        slide_xml
                    ):
                        has_package = True
                        details.append(
                            "OLE package object in: " +
                            slide_file
                        )

                except Exception:
                    pass

            if has_ole:
                findings.append("embedded_ole_object")
            if has_package:
                findings.append("package_object")

    except Exception as e:
        details.append(
            "Embedded object error: " + str(e)
        )

    if not findings:
        details.append(
            "No suspicious embedded objects ✅"
        )

    details.append(
        "Technique 6 findings: " + str(len(findings))
    )
    return findings, details


# ================================================================
# TECHNIQUE 7 — External Link Analysis
# ================================================================

def technique7_external_links(file_path):
    """
    Analyzes all external relationships and links.
    Checks for template injection and malicious URLs.
    """
    findings = []
    details  = []

    details.append("")
    details.append(
        "--- TECHNIQUE 7: EXTERNAL LINK ANALYSIS ---"
    )

    ext = os.path.splitext(file_path)[1].lower()
    if ext not in [
        ".pptx", ".ppsx", ".pptm", ".ppsm", ".potx", ".potm",
    ]:
        details.append(
            "Old format — skipping"
        )
        return findings, details

    try:
        if not zipfile.is_zipfile(file_path):
            details.append("Invalid ZIP — skipping")
            return findings, details

        with zipfile.ZipFile(file_path, "r") as z:
            names   = z.namelist()
            rel_files = [
                n for n in names
                if n.endswith(".rels")
            ]

            details.append(
                "Checking " +
                str(len(rel_files)) +
                " relationship files"
            )

            for rel_file in rel_files:
                try:
                    rel_xml = z.read(
                        rel_file
                    ).decode("utf-8", errors="ignore")

                    # Find all relationship tags
                    rel_tags = re.findall(
                        r'<Relationship\s+[^>]+>',
                        rel_xml,
                        re.IGNORECASE
                    )

                    for rel_tag in rel_tags:
                        if not re.search(r'targetmode=["\']external["\']', rel_tag, re.IGNORECASE):
                            continue

                        target_match = re.search(
                            r'Target=["\']([^"\']+)["\']',
                            rel_tag,
                            re.IGNORECASE
                        )
                        if not target_match:
                            continue

                        url = target_match.group(1).strip()
                        if not url:
                            continue

                        if is_safe_domain(url):
                            continue

                        type_match = re.search(
                            r'Type=["\']([^"\']+)["\']',
                            rel_tag,
                            re.IGNORECASE
                        )
                        rel_type = type_match.group(1).lower() if type_match else ""

                        details.append(
                            "External link: " + url
                        )

                        # Check for template injection in this specific relationship
                        if (
                            "template" in rel_type or
                            "attachedtemplate" in rel_type
                        ):
                            findings.append(
                                "template_injection"
                            )
                            details.append(
                                "CRITICAL: External template "
                                "injection: " + url
                            )
                        elif "hyperlink" not in rel_type:
                            # Non-hyperlink external resource (e.g. oleObject, subdocument, etc.)
                            findings.append(
                                "external_relationship"
                            )

                        # Analyze URL threat
                        url_findings, url_details = (
                            analyze_url(url)
                        )
                        if url_findings:
                            findings.append(
                                "suspicious_external_url"
                            )
                            for d in url_details:
                                details.append(
                                    "  " + d
                                )

                        # Check for image tracker
                        if any(
                            ext in url.lower()
                            for ext in [
                                ".png", ".jpg",
                                ".gif", ".jpeg",
                            ]
                        ):
                            findings.append(
                                "image_tracker"
                            )
                            details.append(
                                "Image tracker URL: " + url
                            )

                except Exception as e:
                    details.append(
                        "Rel file error: " + str(e)
                    )

    except Exception as e:
        details.append(
            "External link error: " + str(e)
        )

    if not findings:
        details.append(
            "No suspicious external links ✅"
        )

    details.append(
        "Technique 7 findings: " + str(len(findings))
    )
    return findings, details


# ================================================================
# TECHNIQUE 8 — Phishing Content Detection
# ================================================================

def technique8_phishing_content(file_path):
    """
    Extracts and analyzes all text from slides
    for phishing content and social engineering.
    """
    findings   = []
    details    = []
    slide_text = ""

    details.append("")
    details.append(
        "--- TECHNIQUE 8: PHISHING CONTENT DETECTION ---"
    )

    # Method 1 — python-pptx for modern formats
    ext = os.path.splitext(file_path)[1].lower()
    if (PPTX_AVAILABLE and
            ext in [
                ".pptx", ".ppsx",
                ".pptm", ".ppsm", ".potx", ".potm",
            ]):
        try:
            prs = Presentation(file_path)

            for slide_num, slide in enumerate(
                prs.slides, 1
            ):
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        slide_text += shape.text_frame.text + " "
                    if shape.has_table:
                        for row in shape.table.rows:
                            for cell in row.cells:
                                slide_text += cell.text_frame.text + " "

            if slide_text.strip():
                details.append(
                    "Text extracted via python-pptx ✅"
                )

        except Exception as e:
            details.append(
                "python-pptx error: " + str(e)
            )

    # Method 2 — ZIP XML extraction fallback or supplement
    if not slide_text:
        try:
            if zipfile.is_zipfile(file_path):
                with zipfile.ZipFile(
                    file_path, "r"
                ) as z:
                    for name in z.namelist():
                        if re.match(
                            r'ppt/slides/slide\d+\.xml',
                            name
                        ):
                            xml = z.read(name).decode(
                                "utf-8", errors="ignore"
                            )
                            text_matches = re.findall(
                                r'<a:t>([^<]+)</a:t>',
                                xml
                            )
                            slide_text += " ".join(
                                text_matches
                            ) + " "

                if slide_text.strip():
                    details.append(
                        "Text extracted via XML ✅"
                    )

        except Exception as e:
            details.append(
                "XML extraction error: " + str(e)
            )

    if not slide_text.strip():
        details.append(
            "No text found in presentation"
        )
        details.append(
            "Technique 8 findings: " + str(len(findings))
        )
        return findings, details

    word_count = len(slide_text.split())
    details.append(
        "Total words in slides: " + str(word_count)
    )

    text_lower     = slide_text.lower()
    urgency_count      = 0
    finance_count      = 0
    cred_count         = 0
    enable_macro_count = 0

    for category, keywords in PHISHING_KEYWORDS.items():
        for keyword in keywords:
            count = text_lower.count(keyword)
            if count > 0:
                findings.append("phishing_keyword")
                details.append(
                    "[" + category + "] '" +
                    keyword + "' x" + str(count)
                )
                if category == "urgency_phrases":
                    urgency_count += count
                elif category == "financial_terms":
                    finance_count += count
                elif category == "credential_harvesting":
                    cred_count += count
                elif category == "macro_lure":
                    enable_macro_count += count

    if urgency_count >= 2:
        findings.append("urgent_tone_detected")
        details.append(
            "Urgency tone: " + str(urgency_count) +
            " phrases"
        )

    if finance_count >= 2:
        findings.append("financial_terms_detected")
        details.append(
            "Financial targeting: " +
            str(finance_count) + " terms"
        )

    if cred_count >= 1:
        findings.append("credential_harvesting")
        details.append(
            "Credential harvesting: " +
            str(cred_count) + " phrases"
        )

    if enable_macro_count >= 1:
        findings.append("enable_macro_lure")
        details.append(
            "CRITICAL: Macro lure text found (" +
            str(enable_macro_count) + " phrases) — tricks user into enabling macros!"
        )

    if not findings:
        details.append(
            "No phishing content detected ✅"
        )

    details.append(
        "Technique 8 findings: " + str(len(findings))
    )
    return findings, details


# ================================================================
# TECHNIQUE 9 — URL Detection in Slides
# ================================================================

def technique9_url_detection(file_path):
    """
    Finds and analyzes all URLs in slide content and slide relationships.
    """
    findings = []
    details  = []

    details.append("")
    details.append(
        "--- TECHNIQUE 9: URL DETECTION IN SLIDES ---"
    )

    all_urls = []

    # Method 1 — python-pptx hyperlinks
    ext = os.path.splitext(file_path)[1].lower()
    if (PPTX_AVAILABLE and
            ext in [
                ".pptx", ".ppsx",
                ".pptm", ".ppsm", ".potx", ".potm",
            ]):
        try:
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in (
                            shape.text_frame.paragraphs
                        ):
                            for run in para.runs:
                                if (
                                    run.hyperlink and
                                    run.hyperlink.address
                                ):
                                    addr = run.hyperlink.address.strip()
                                    if addr and addr not in all_urls:
                                        all_urls.append(addr)
                    if shape.has_table:
                        for row in shape.table.rows:
                            for cell in row.cells:
                                for para in cell.text_frame.paragraphs:
                                    for run in para.runs:
                                        if (
                                            run.hyperlink and
                                            run.hyperlink.address
                                        ):
                                            addr = run.hyperlink.address.strip()
                                            if addr and addr not in all_urls:
                                                all_urls.append(addr)

        except Exception as e:
            details.append(
                "python-pptx URL error: " + str(e)
            )

    # Method 2 — URL extraction from slide rels and slide XML
    try:
        if zipfile.is_zipfile(file_path):
            with zipfile.ZipFile(file_path, "r") as z:
                for name in z.namelist():
                    if (
                        re.match(r'ppt/slides/slide\d+\.xml', name) or
                        re.match(r'ppt/slides/_rels/slide\d+\.xml\.rels', name)
                    ):
                        xml = z.read(name).decode(
                            "utf-8", errors="ignore"
                        )
                        url_pattern = re.compile(
                            r'https?://[^\s<>"{}|\\^`\[\]]+'
                        )
                        for url in url_pattern.findall(xml):
                            # Skip placeholder variables and localhost
                            if "$" in url or "localhost" in url.lower():
                                continue
                            if url not in all_urls:
                                all_urls.append(url)

    except Exception as e:
        details.append("XML URL error: " + str(e))

    details.append(
        "Total URLs found: " + str(len(all_urls))
    )

    for url in all_urls:
        if is_safe_domain(url):
            details.append(
                "Safe URL: " + url[:60]
            )
            continue

        url_findings, url_details = analyze_url(url)
        if url_findings:
            findings.extend(url_findings)
            for d in url_details:
                details.append("  " + d)
        else:
            details.append(
                "URL (clean): " + url[:60]
            )

    if not all_urls:
        details.append(
            "No URLs found in slides ✅"
        )

    details.append(
        "Technique 9 findings: " + str(len(findings))
    )
    return findings, details


# ================================================================
# TECHNIQUE 10 — Hidden Slide Detection (PPT UNIQUE)
# ================================================================

def technique10_hidden_slides(file_path):
    """
    Detects hidden slides in presentation.
    This technique is UNIQUE to PowerPoint files!

    Attackers hide malicious content in hidden slides.
    Victim sees normal presentation.
    Macros access hidden slide content secretly.
    """
    findings = []
    details  = []

    details.append("")
    details.append(
        "--- TECHNIQUE 10: HIDDEN SLIDE DETECTION ---"
    )

    ext = os.path.splitext(file_path)[1].lower()
    if ext not in [
        ".pptx", ".ppsx", ".pptm", ".ppsm", ".potx", ".potm",
    ]:
        details.append(
            "Old format — hidden slide check skipped"
        )
        return findings, details

    try:
        if not zipfile.is_zipfile(file_path):
            details.append("Invalid ZIP — skipping")
            return findings, details

        # Check presentation.xml for hidden slides
        with zipfile.ZipFile(file_path, "r") as z:
            names = z.namelist()

            if "ppt/presentation.xml" not in names:
                details.append(
                    "presentation.xml not found"
                )
                return findings, details

            prs_xml = z.read(
                "ppt/presentation.xml"
            ).decode("utf-8", errors="ignore")

            # Find all sldId entries
            # show="0" or show="false" means hidden slide
            hidden_pattern = re.compile(
                r'<p:sldId[^>]*show="(?:0|false)"[^>]*\/?>',
                re.IGNORECASE
            )
            hidden_slides = hidden_pattern.findall(
                prs_xml
            )

            # Count total slides
            total_pattern = re.compile(
                r'<p:sldId[^>]*\/?>'
            )
            total_slides = total_pattern.findall(
                prs_xml
            )

            details.append(
                "Total slides  : " +
                str(len(total_slides))
            )
            details.append(
                "Hidden slides : " +
                str(len(hidden_slides))
            )

            if hidden_slides:
                findings.append("hidden_slide")
                details.append(
                    str(len(hidden_slides)) +
                    " hidden slide(s) detected!"
                )

                # Map r:id to actual slide file path via presentation.xml.rels
                rid_to_slide = {}
                if "ppt/_rels/presentation.xml.rels" in names:
                    try:
                        rels_xml = z.read("ppt/_rels/presentation.xml.rels").decode("utf-8", errors="ignore")
                        for rel in re.finditer(r'<Relationship\s+[^>]*Id="([^"]+)"[^>]*Target="([^"]+)"', rels_xml, re.IGNORECASE):
                            rid, target = rel.group(1), rel.group(2)
                            target_norm = target.lstrip("/")
                            if not target_norm.startswith("ppt/"):
                                target_norm = "ppt/" + target_norm
                            rid_to_slide[rid] = target_norm
                    except Exception:
                        pass

                hidden_slide_files = []
                for sld_tag in hidden_slides:
                    rid_match = re.search(r'r:id="([^"]+)"', sld_tag, re.IGNORECASE)
                    if rid_match and rid_match.group(1) in rid_to_slide:
                        hidden_slide_files.append(rid_to_slide[rid_match.group(1)])

                slide_files = [
                    n for n in names
                    if re.match(
                        r'ppt/slides/slide\d+\.xml', n
                    )
                ]

                # Fallback to all slide files if rels resolution was unavailable
                if not hidden_slide_files:
                    hidden_slide_files = slide_files

                for slide_file in hidden_slide_files:
                    if slide_file in names:
                        try:
                            slide_xml = z.read(
                                slide_file
                            ).decode(
                                "utf-8", errors="ignore"
                            )
                            text_matches = re.findall(
                                r'<a:t>([^<]+)</a:t>',
                                slide_xml
                            )
                            text = " ".join(
                                text_matches
                            ).strip()

                            if text:
                                findings.append(
                                    "hidden_slide_with_content"
                                )
                                details.append(
                                    "CRITICAL: Hidden slide "
                                    "has content: " +
                                    text[:100]
                                )
                                break

                        except Exception:
                            pass
            else:
                details.append(
                    "No hidden slides found ✅"
                )

    except Exception as e:
        details.append(
            "Hidden slide error: " + str(e)
        )

    details.append(
        "Technique 10 findings: " + str(len(findings))
    )
    return findings, details


# ================================================================
# TECHNIQUE 11 — Obfuscation Detection
# ================================================================

def technique11_obfuscation(file_path, vba_content):
    """
    Detects code obfuscation techniques.
    Base64, character concat, string splitting.
    """
    findings = []
    details  = []

    details.append("")
    details.append(
        "--- TECHNIQUE 11: OBFUSCATION DETECTION ---"
    )

    content_to_check = vba_content or ""

    # Check slide text content
    ext = os.path.splitext(file_path)[1].lower()
    try:
        if (ext in [
            ".pptx", ".ppsx", ".pptm", ".ppsm", ".potx", ".potm",
        ] and zipfile.is_zipfile(file_path)):
            with zipfile.ZipFile(
                file_path, "r"
            ) as z:
                for name in z.namelist():
                    if re.match(r'ppt/slides/slide\d+\.xml', name):
                        try:
                            content_to_check += z.read(
                                name
                            ).decode(
                                "utf-8", errors="ignore"
                            )
                        except Exception:
                            pass
    except Exception:
        pass

    if not content_to_check:
        details.append("No content to analyze")
        details.append(
            "Technique 11 findings: " +
            str(len(findings))
        )
        return findings, details

    # Check 1 — Base64 strings (in VBA or slide text)
    b64_pattern = re.compile(
        r'[A-Za-z0-9+/]{60,}={0,2}'
    )
    b64_matches = b64_pattern.findall(content_to_check)
    if b64_matches:
        findings.append("base64_payload")
        details.append(
            "Base64 encoded strings found: " +
            str(len(b64_matches))
        )
        details.append(
            "Sample: " + b64_matches[0][:50]
        )

    # Check 2 — Chr() concatenation
    chr_pattern = re.compile(
        r'chr[w\$]?\(\s*\d+\s*\)', re.IGNORECASE
    )
    chr_matches = chr_pattern.findall(
        content_to_check
    )
    if len(chr_matches) > 5:
        findings.append("char_concat_obfuscation")
        details.append(
            "Chr() concatenation: " +
            str(len(chr_matches)) + " instances"
        )

    # Check 3 — String splitting
    split_patterns = [
        r'"[a-zA-Z]{1,4}"\s*&\s*"[a-zA-Z]{1,4}"',
        r'"pow"\s*&\s*"er"',
        r'"cmd"\s*&\s*"\.exe"',
    ]
    for pattern in split_patterns:
        if re.search(
            pattern,
            content_to_check,
            re.IGNORECASE
        ):
            findings.append(
                "string_split_obfuscation"
            )
            details.append(
                "String splitting obfuscation detected!"
            )
            break

    # Check 4 — High entropy strings in VBA
    if vba_content:
        lines = vba_content.split("\n")
        for line in lines:
            if len(line.strip()) > 50:
                entropy = calculate_entropy(
                    line.strip().encode()
                )
                if entropy > 5.5:
                    findings.append(
                        "high_entropy_string"
                    )
                    details.append(
                        "High entropy string: " +
                        line[:50].strip() + "..."
                    )
                    break

    if not findings:
        details.append(
            "No obfuscation detected ✅"
        )

    details.append(
        "Technique 11 findings: " + str(len(findings))
    )
    return findings, details


# ================================================================
# TECHNIQUE 12 — Action Button Analysis (PPT UNIQUE)
# ================================================================

def technique12_action_buttons(file_path):
    """
    Detects malicious action buttons in slides.
    This technique is UNIQUE to PowerPoint files!

    Common attack:
    Invisible button covers slide.
    Victim clicks anywhere.
    Button runs malware or macro silently.
    """
    findings = []
    details  = []

    details.append("")
    details.append(
        "--- TECHNIQUE 12: ACTION BUTTON ANALYSIS ---"
    )

    ext = os.path.splitext(file_path)[1].lower()
    if ext not in [
        ".pptx", ".ppsx", ".pptm", ".ppsm", ".potx", ".potm",
    ]:
        details.append(
            "Old format — action button check skipped"
        )
        return findings, details

    try:
        if not zipfile.is_zipfile(file_path):
            details.append("Invalid ZIP — skipping")
            return findings, details

        with zipfile.ZipFile(file_path, "r") as z:
            names = z.namelist()

            slide_files = [
                n for n in names
                if re.match(
                    r'ppt/slides/slide\d+\.xml', n
                )
            ]

            for slide_file in slide_files:
                try:
                    slide_xml = z.read(
                        slide_file
                    ).decode("utf-8", errors="ignore")

                    slide_num = re.search(
                        r'slide(\d+)\.xml', slide_file
                    )
                    snum = (
                        slide_num.group(1)
                        if slide_num else "?"
                    )

                    # Check for shape click action (on shape non-visual properties or explicit action)
                    has_shape_click = bool(
                        re.search(r'<p:cNv(?:Sp)?Pr\b[^>]*>.*?<a:hlinkClick\b', slide_xml, re.DOTALL | re.IGNORECASE) or
                        re.search(r'<a:hlinkClick\b[^>]*action=', slide_xml, re.IGNORECASE)
                    )
                    if has_shape_click:
                        findings.append(
                            "action_button_found"
                        )
                        details.append(
                            "Slide " + snum +
                            ": click action button found on shape"
                        )

                    # Check for shape mouseover action
                    has_shape_hover = bool(
                        re.search(r'<p:cNv(?:Sp)?Pr\b[^>]*>.*?<a:hlinkMouseOver\b', slide_xml, re.DOTALL | re.IGNORECASE) or
                        re.search(r'<a:hlinkMouseOver\b[^>]*action=', slide_xml, re.IGNORECASE)
                    )
                    if has_shape_hover:
                        findings.append(
                            "mouseover_action"
                        )
                        details.append(
                            "CRITICAL: Slide " + snum +
                            ": mouseover action found!"
                        )

                    # Check for run macro action
                    if "ppaction://macro" in slide_xml.lower():
                        findings.append(
                            "macro_action_button"
                        )
                        details.append(
                            "CRITICAL: Slide " + snum +
                            ": macro trigger button!"
                        )

                    # Check for run program action (Note: hlinksldjump is benign slide jump)
                    if "ppaction://program" in slide_xml.lower():
                        findings.append(
                            "run_program_action"
                        )
                        details.append(
                            "CRITICAL: Slide " + snum +
                            ": run program action!"
                        )

                    # Check for invisible shapes with actions
                    # Extract shape blocks <p:sp>...</p:sp>
                    shape_blocks = re.findall(
                        r'<p:sp\b[^>]*>.*?</p:sp>',
                        slide_xml,
                        re.DOTALL
                    )
                    for shape_xml in shape_blocks:
                        # Has action?
                        shape_has_action = (
                            "hlinkclick" in shape_xml.lower() or
                            "hlinkmouseover" in shape_xml.lower() or
                            "ppaction://" in shape_xml.lower()
                        )
                        if not shape_has_action:
                            continue

                        # Check if shape has text (visible text = normal button, not invisible trap)
                        text_matches = re.findall(
                            r'<a:t>([^<]+)</a:t>',
                            shape_xml
                        )
                        has_text = bool("".join(text_matches).strip())

                        # Check for noFill / transparent fill
                        has_no_fill = "<a:noFill/>" in shape_xml or "<a:noFill " in shape_xml

                        if has_no_fill and not has_text:
                            findings.append(
                                "invisible_button"
                            )
                            details.append(
                                "Slide " + snum +
                                ": invisible button WITH action (no text, no fill)!"
                            )
                            break

                except Exception as e:
                    details.append(
                        "Slide error: " + str(e)
                    )

    except Exception as e:
        details.append(
            "Action button error: " + str(e)
        )

    if not findings:
        details.append(
            "No suspicious action buttons ✅"
        )

    details.append(
        "Technique 12 findings: " + str(len(findings))
    )
    return findings, details


# ================================================================
# TECHNIQUE 13 — Media File Analysis
# ================================================================

def technique13_media_files(file_path):
    """
    Analyzes media files embedded in presentation.
    Detects disguised executables and high-entropy media.
    """
    findings = []
    details  = []

    details.append("")
    details.append(
        "--- TECHNIQUE 13: MEDIA FILE ANALYSIS ---"
    )

    ext = os.path.splitext(file_path)[1].lower()
    if ext not in [
        ".pptx", ".ppsx", ".pptm", ".ppsm", ".potx", ".potm",
    ]:
        details.append(
            "Old format — media check skipped"
        )
        return findings, details

    try:
        if not zipfile.is_zipfile(file_path):
            details.append("Invalid ZIP — skipping")
            return findings, details

        with zipfile.ZipFile(file_path, "r") as z:
            names = z.namelist()

            media_files = [
                n for n in names
                if "ppt/media/" in n
            ]

            details.append(
                "Media files found: " +
                str(len(media_files))
            )

            for media in media_files:
                media_ext = os.path.splitext(
                    media
                )[1].lower()
                media_name = os.path.basename(media)

                try:
                    content = z.read(media)
                    size_kb = round(
                        len(content) / 1024, 2
                    )
                    details.append(
                        "Media: " + media_name +
                        " (" + str(size_kb) + " KB)"
                    )

                    # Check file signature vs extension
                    if (
                        media_ext in [
                            ".png", ".jpg",
                            ".jpeg", ".gif"
                        ]
                    ):
                        # JPG signature
                        if (
                            media_ext in [
                                ".jpg", ".jpeg"
                            ] and
                            content[:3] != b"\xff\xd8\xff"
                        ):
                            findings.append(
                                "media_type_mismatch"
                            )
                            details.append(
                                "CRITICAL: " +
                                media_name +
                                " fails JPG signature!"
                            )

                        # PNG signature
                        elif (
                            media_ext == ".png" and
                            content[:4] != b"\x89PNG"
                        ):
                            findings.append(
                                "media_type_mismatch"
                            )
                            details.append(
                                "CRITICAL: " +
                                media_name +
                                " fails PNG signature!"
                            )

                    # Check for MZ header in media
                    if content[:2] == b"MZ":
                        findings.append(
                            "embedded_executable"
                        )
                        findings.append(
                            "mz_header_found"
                        )
                        details.append(
                            "CRITICAL: EXE header in " +
                            media_name + "!"
                        )

                    # Check for large media files (> 10MB)
                    if size_kb > 10240:
                        findings.append(
                            "large_media_file"
                        )
                        details.append(
                            "Large media: " +
                            media_name +
                            " (" + str(size_kb) + " KB)"
                        )

                    # Check entropy for non-standard or suspicious media
                    if media_ext not in [".png", ".jpg", ".jpeg", ".mp4", ".mp3", ".webp"]:
                        entropy = calculate_entropy(
                            content[:2000]
                        )
                        if entropy > 7.9:
                            findings.append(
                                "high_entropy_media"
                            )
                            details.append(
                                "High entropy media: " +
                                media_name +
                                " (" + str(entropy) + ")"
                            )

                except Exception as e:
                    details.append(
                        "Media read error: " + str(e)
                    )

    except Exception as e:
        details.append(
            "Media analysis error: " + str(e)
        )

    if not findings:
        details.append(
            "No suspicious media files ✅"
        )

    details.append(
        "Technique 13 findings: " + str(len(findings))
    )
    return findings, details


# ================================================================
# TECHNIQUE 14 — Attack Chain Inference + Scoring
# ================================================================

def technique14_attack_chain(all_findings):
    """
    Infers complete attack patterns from all findings.
    """
    findings = []
    details  = []

    details.append("")
    details.append(
        "--- TECHNIQUE 14: ATTACK CHAIN INFERENCE ---"
    )

    fs = set(all_findings)

    # Chain 1 — Dropper attack
    if (
        "vba_macro_detected" in fs and
        ("shell_command" in fs or "process_creation" in fs or "powershell_in_vba" in fs)
    ):
        findings.append("dropper_pattern")
        details.append(
            "ATTACK CHAIN: Macro present -> "
            "Shell/Process execution -> Dropper attack!"
        )

    # Chain 2 — Remote template injection
    if "template_injection" in fs:
        findings.append("remote_template_attack")
        details.append(
            "ATTACK CHAIN: External template -> "
            "Loads remote payload -> "
            "Remote template attack!"
        )

    # Chain 3 — Social engineering
    if (
        "credential_harvesting" in fs and
        "phishing_keyword" in fs
    ):
        findings.append("social_engineering_chain")
        details.append(
            "ATTACK CHAIN: Phishing text -> "
            "Credential lure -> "
            "Social engineering attack!"
        )

    # Chain 4 — Hidden payload
    if (
        "hidden_slide" in fs and
        "vba_macro_detected" in fs
    ):
        findings.append("hidden_payload_chain")
        details.append(
            "ATTACK CHAIN: Hidden slide -> "
            "Macro accesses it -> "
            "Hidden payload attack!"
        )

    # Chain 5 — Click to execute
    if (
        "macro_action_button" in fs or
        "run_program_action" in fs or
        ("action_button_found" in fs and (
            "vba_macro_detected" in fs or
            "suspicious_url" in fs or
            "suspicious_external_url" in fs or
            "embedded_executable" in fs or
            "mz_header_found" in fs
        ))
    ):
        findings.append("click_execute_chain")
        details.append(
            "ATTACK CHAIN: Action button -> "
            "Victim interaction -> "
            "Malware executes!"
        )

    # Chain 6 — Multistage
    if (
        "base64_payload" in fs and
        (
            "network_call_in_vba" in fs or
            "suspicious_external_url" in fs or
            "template_injection" in fs
        )
    ):
        findings.append("multistage_indicator")
        details.append(
            "ATTACK CHAIN: Encoded payload -> "
            "Network download -> "
            "Multistage attack!"
        )

    if not findings:
        details.append(
            "No clear attack chain identified"
        )

    details.append(
        "Technique 14 findings: " + str(len(findings))
    )
    return findings, details


# ================================================================
# SCORING FUNCTION
# ================================================================

# DEPRECATED: calculate_final_score
# NOTE: Scoring is now handled by scorer.py (centralized system)
def calculate_final_score(all_findings):
    return 0, "Deprecated", "LOW", {}


# ================================================================
# MAIN FUNCTION — parse_ppt
# Called by app.py when user uploads PowerPoint file
# ================================================================

def parse_ppt(file_path):
    """
    file_path = full path to uploaded PowerPoint file
    Returns   = complete analysis result dict
    """

    all_findings = []
    all_details  = []
    sha256_hash  = ""

    if not file_path or not os.path.exists(file_path):
        return {
            "findings": ["invalid_ppt_format"],
            "details": ["File not found or invalid path: " + str(file_path)],
            "sha256": "",
        }

    try:

        all_details.append("=" * 55)
        all_details.append(
            "DARKHOOK_ DEFENCE — PPT FILE ANALYSIS"
        )
        all_details.append("=" * 55)
        all_details.append(
            "File: " + os.path.basename(file_path)
        )

        with open(file_path, "rb") as f:
            sha256_hash = hashlib.sha256(
                f.read()
            ).hexdigest()

        all_details.append(
            "SHA256: " + sha256_hash
        )
        all_details.append("")

        # Technique 1 — File validation
        f1, d1 = technique1_file_validation(file_path)
        all_findings.extend(f1)
        all_details.extend(d1)

        # Technique 2 — Metadata
        f2, d2 = technique2_metadata(file_path)
        all_findings.extend(f2)
        all_details.extend(d2)

        # Technique 3 — Macro detection
        # Returns vba_content for other techniques
        f3, d3, vba_content = (
            technique3_macro_detection(file_path)
        )
        all_findings.extend(f3)
        all_details.extend(d3)

        # Technique 4 — VBA behavior
        f4, d4 = technique4_vba_behavior(vba_content)
        all_findings.extend(f4)
        all_details.extend(d4)

        # Technique 5 — Animation triggers (PPT UNIQUE)
        f5, d5 = technique5_animation_triggers(
            file_path
        )
        all_findings.extend(f5)
        all_details.extend(d5)

        # Technique 6 — Embedded objects
        f6, d6 = technique6_embedded_objects(file_path)
        all_findings.extend(f6)
        all_details.extend(d6)

        # Technique 7 — External links
        f7, d7 = technique7_external_links(file_path)
        all_findings.extend(f7)
        all_details.extend(d7)

        # Technique 8 — Phishing content
        f8, d8 = technique8_phishing_content(file_path)
        all_findings.extend(f8)
        all_details.extend(d8)

        # Technique 9 — URL detection
        f9, d9 = technique9_url_detection(file_path)
        all_findings.extend(f9)
        all_details.extend(d9)

        # Technique 10 — Hidden slides (PPT UNIQUE)
        f10, d10 = technique10_hidden_slides(file_path)
        all_findings.extend(f10)
        all_details.extend(d10)

        # Technique 11 — Obfuscation
        f11, d11 = technique11_obfuscation(
            file_path, vba_content
        )
        all_findings.extend(f11)
        all_details.extend(d11)

        # Technique 12 — Action buttons (PPT UNIQUE)
        f12, d12 = technique12_action_buttons(
            file_path
        )
        all_findings.extend(f12)
        all_details.extend(d12)

        # Technique 13 — Media files
        f13, d13 = technique13_media_files(file_path)
        all_findings.extend(f13)
        all_details.extend(d13)

        # Technique 14 — Attack chain
        f14, d14 = technique14_attack_chain(
            all_findings
        )
        all_findings.extend(f14)
        all_details.extend(d14)

        all_details.append("")
        all_details.append("--- ANALYSIS COMPLETE ---")
        all_details.append("Total techniques run : 14")
        all_details.append("Total findings       : " + str(len(all_findings)))
        all_details.append("")
        all_details.append("Note: Final scoring calculated by centralized scorer.py")

    except Exception as error:
        all_details.append(
            "Critical error: " + str(error)
        )

    return {
        "findings" : all_findings,
        "details"  : all_details,
        "sha256"   : sha256_hash,
    }