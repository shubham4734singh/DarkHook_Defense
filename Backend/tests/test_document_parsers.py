import os
import sys
import tempfile
import pytest
import openpyxl
import fitz
import docx
from pptx import Presentation
from PIL import Image, ImageDraw

# Add Backend to path
sys.path.insert(0, os.path.abspath(os.path.join(os.path.dirname(__file__), "..")))

from services.document_parsers.scorer import calculate_score, get_verdict, WEIGHTS
from services.document_parsers.pdf_parser import parse_pdf
from services.document_parsers.docx_parser import parse_docx
from services.document_parsers.excel_parser import parse_excel
from services.document_parsers.ppt_parser import parse_ppt
from services.document_parsers.ocr_parser import parse_image
from services.document_analyzer import document_analyzer


def test_scorer_verdicts():
    assert get_verdict(0) == "Safe"
    assert get_verdict(39) == "Safe"
    assert get_verdict(40) == "Suspicious"
    assert get_verdict(69) == "Suspicious"
    assert get_verdict(70) == "Phishing"
    assert get_verdict(100) == "Phishing"

    res = calculate_score(["suspicious_url", "malicious_macro"])
    assert res["score"] == 55
    assert res["verdict"] == "Suspicious"
    assert res["severity"] == "MEDIUM"
    assert "suspicious_url" in res["breakdown"]
    assert "malicious_macro" in res["breakdown"]


def test_pdf_parser():
    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((50, 50), "Important account verification required immediately.")
    with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as f:
        pdf_path = f.name
    doc.save(pdf_path)
    doc.close()

    try:
        res = parse_pdf(pdf_path)
        assert isinstance(res, dict)
        assert "findings" in res
        assert "details" in res
        assert "sha256" in res
        assert len(res["sha256"]) == 64
    finally:
        if os.path.exists(pdf_path):
            os.unlink(pdf_path)


def test_docx_parser():
    doc = docx.Document()
    doc.add_paragraph("Dear user, please verify your account and enter your credentials.")
    with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as f:
        docx_path = f.name
    doc.save(docx_path)

    try:
        res = parse_docx(docx_path)
        assert isinstance(res, dict)
        assert "findings" in res
        assert "details" in res
        assert "sha256" in res
        assert len(res["sha256"]) == 64
        assert "phishing_keyword" in res["findings"]
    finally:
        if os.path.exists(docx_path):
            os.unlink(docx_path)


def test_excel_parser_hyperlink_and_cleanup():
    wb = openpyxl.Workbook()
    ws = wb.active
    ws['A1'] = "Test Document"
    ws['A2'] = "=HYPERLINK(\"http://phishing-site.xyz/login\", \"Click Here to Update Password\")"
    ws['A3'] = "account suspended"
    
    with tempfile.NamedTemporaryFile(suffix=".xlsx", delete=False) as f:
        xlsx_path = f.name
    wb.save(xlsx_path)
    wb.close()

    try:
        res = parse_excel(xlsx_path)
        assert isinstance(res, dict)
        assert "findings" in res
        assert "details" in res
        assert "sha256" in res
        assert len(res["sha256"]) == 64
        assert "formula_hyperlink_injection" in res["findings"] or "suspicious_url" in res["findings"]
    finally:
        # File must not be locked by openpyxl
        if os.path.exists(xlsx_path):
            os.unlink(xlsx_path)


def test_ppt_parser():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Urgent: verify your account"
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        pptx_path = f.name
    prs.save(pptx_path)

    try:
        res = parse_ppt(pptx_path)
        assert isinstance(res, dict)
        assert "findings" in res
        assert "details" in res
        assert "sha256" in res
        assert len(res["sha256"]) == 64
        assert "phishing_keyword" in res["findings"]

        # Test non-existent file handling
        missing_res = parse_ppt("non_existent_file.pptx")
        assert missing_res["sha256"] == ""
        assert "invalid_ppt_format" in missing_res["findings"]
    finally:
        if os.path.exists(pptx_path):
            os.unlink(pptx_path)


def test_ppt_parser_slide_jump_and_actions():
    import zipfile
    import io
    prs = Presentation()
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Table of Contents"

    bio = io.BytesIO()
    prs.save(bio)
    bio.seek(0)

    with zipfile.ZipFile(bio, "r") as zin:
        file_map = {item.filename: zin.read(item.filename) for item in zin.infolist()}

    s1_xml = file_map["ppt/slides/slide1.xml"].decode("utf-8")
    s1_xml = s1_xml.replace("</p:sp>", '<a:hlinkClick action="ppaction://hlinksldjump?slideindex=2"/></p:sp>')
    file_map["ppt/slides/slide1.xml"] = s1_xml.encode("utf-8")

    out_bio = io.BytesIO()
    with zipfile.ZipFile(out_bio, "w", zipfile.ZIP_DEFLATED) as zout:
        for fname, data in file_map.items():
            zout.writestr(fname, data)

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        path = f.name
        f.write(out_bio.getvalue())

    try:
        res = parse_ppt(path)
        # Benign slide jump must NOT be flagged as run program
        assert "run_program_action" not in res["findings"]
    finally:
        if os.path.exists(path):
            os.unlink(path)


def test_ppt_parser_template_injection():
    import zipfile
    import io
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[0])

    bio = io.BytesIO()
    prs.save(bio)
    bio.seek(0)

    with zipfile.ZipFile(bio, "r") as zin:
        file_map = {item.filename: zin.read(item.filename) for item in zin.infolist()}

    # Add external template relationship
    rel_content = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>\n'
        '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
        '<Relationship Id="rId99" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/attachedTemplate" '
        'Target="http://malicious-site.xyz/template.dotm" TargetMode="External"/>'
        '</Relationships>'
    )
    file_map["ppt/_rels/presentation.xml.rels"] = rel_content.encode("utf-8")

    out_bio = io.BytesIO()
    with zipfile.ZipFile(out_bio, "w", zipfile.ZIP_DEFLATED) as zout:
        for fname, data in file_map.items():
            zout.writestr(fname, data)

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        path = f.name
        f.write(out_bio.getvalue())

    try:
        res = parse_ppt(path)
        assert "template_injection" in res["findings"]
        assert "remote_template_attack" in res["findings"]
    finally:
        if os.path.exists(path):
            os.unlink(path)



def test_ocr_parser():
    img = Image.new("RGB", (300, 100), color=(255, 255, 255))
    d = ImageDraw.Draw(img)
    d.text((10, 10), "Account Suspended", fill=(0, 0, 0))
    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as f:
        img_path = f.name
    img.save(img_path)

    try:
        res = parse_image(img_path)
        assert isinstance(res, dict)
        assert "findings" in res
        assert "details" in res
        assert "sha256" in res
        assert len(res["sha256"]) == 64
    finally:
        if os.path.exists(img_path):
            os.unlink(img_path)


def test_document_analyzer_integration():
    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((50, 50), "Hello normal clean document with link http://phishing-update.xyz/login.")
    data = doc.tobytes()
    doc.close()

    result = document_analyzer.analyze_document("test.pdf", data)
    assert result["fileName"] == "test.pdf"
    assert "riskScore" in result
    assert "verdict" in result
    assert "findings" in result
    assert "mitreTechniques" in result
    assert "extractedUrls" in result
    assert len(result["extractedUrls"]) > 0
    assert result["extractedUrls"][0]["is_suspicious"] is True


def test_magic_byte_masquerading_detection():
    # Disguised Windows PE binary (MZ header) passed as .pdf
    fake_pdf = b"MZ\x90\x00\x03\x00\x00\x00\x04\x00\x00\x00\xff\xff\x00\x00" + b"\x00" * 100
    res = document_analyzer.analyze_document("invoice.pdf", fake_pdf)
    assert "file_type_mismatch" in res["findings"]
    assert "embedded_executable" in res["findings"]
    assert res["riskScore"] >= 70
    assert res["verdict"] == "Phishing"

